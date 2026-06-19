import os
import tempfile

import pytest

from app import extract


def _tmp(suffix: str) -> str:
    f = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    f.close()
    return f.name


def test_txt():
    p = _tmp(".txt")
    with open(p, "w", encoding="utf-8") as f:
        f.write("Xin chào thế giới")
    try:
        assert "thế giới" in extract.extract_text(p, "a.txt")
    finally:
        os.unlink(p)


def test_docx():
    import docx

    p = _tmp(".docx")
    d = docx.Document()
    d.add_paragraph("Nội dung tờ trình ABC")
    d.save(p)
    try:
        assert "tờ trình ABC" in extract.extract_text(p, "a.docx")
    finally:
        os.unlink(p)


def test_xlsx():
    from openpyxl import Workbook

    p = _tmp(".xlsx")
    wb = Workbook()
    ws = wb.active
    ws.append(["Chi phí", 1000])
    wb.save(p)
    try:
        assert "Chi phí" in extract.extract_text(p, "a.xlsx")
    finally:
        os.unlink(p)


def test_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    p = _tmp(".pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "Slide nội dung XYZ"
    prs.save(p)
    try:
        assert "XYZ" in extract.extract_text(p, "a.pptx")
    finally:
        os.unlink(p)


def test_reject_bad_ext():
    with pytest.raises(ValueError):
        extract.extract_text("/x/y.exe", "y.exe")
