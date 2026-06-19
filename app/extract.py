"""Trích text từ nhiều loại tài liệu (PDF/DOCX/PPTX/XLSX/TXT/MD)."""
import os

ALLOWED_DOC_EXT = {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md"}


def extract_text(file_path: str, filename: str) -> str:
    ext = os.path.splitext(filename or "")[1].lower()
    if ext == ".pdf":
        return _pdf(file_path)
    if ext == ".docx":
        return _docx(file_path)
    if ext == ".pptx":
        return _pptx(file_path)
    if ext == ".xlsx":
        return _xlsx(file_path)
    if ext in (".txt", ".md"):
        return _txt(file_path)
    raise ValueError(f"Định dạng không hỗ trợ: {ext or '(không rõ)'}")


def _pdf(path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    return "\n".join((page.extract_text() or "") for page in reader.pages).strip()


def _docx(path: str) -> str:
    import docx

    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs]
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts).strip()


def _xlsx(path: str) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"--- Sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts).strip()


def _txt(path: str) -> str:
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read().strip()
